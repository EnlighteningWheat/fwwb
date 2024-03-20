from pptx import Presentation
import json
import os



def pptf_to_json(folder_path,out_path):
    # 获取文件夹下所有文件名
    files_in_folder = os.listdir(folder_path)
    folder_name = os.path.basename(folder_path)
    # 遍历并输出每个文件的完整路径
    for file in files_in_folder:
        ppt_path = os.path.join(folder_path, file)
        json_filename = os.path.splitext(file)[0] + ".json" # 替换为你想要输出的JSON文件路径
        json_path=os.path.join(out_path, json_filename)
        ppt_to_json(ppt_path,json_path)
def ppt_to_json(ppt_path, json_path):
    """
    将PPT中的文本内容提取并保存为JSON格式。

    参数:
    ppt_path: str
        PPT文件的路径。
    json_path: str
        输出的JSON文件的路径。
    """
    # 读取PPT文件
    prs = Presentation(ppt_path)

    # 提取文本内容
    slides_text = []
    for slide_number, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            slide_content = {"text":''}
            if hasattr(shape, "text") and shape.text.strip()!='' and len(shape.text.strip())>4:
                slide_content["text"]=(shape.text.replace('\n', '').replace(' ', ''))
                slides_text.append(slide_content)

    # 转换为JSON并写入文件
    with open(json_path, "w", encoding="utf-8") as json_file:
        json.dump(slides_text, json_file, ensure_ascii=False, indent=4)


if __name__ == '__main__':
# 使用示例
#     ppt_path = "第一讲：有用之用与无用之用.pptx"  # 替换为你的PPT文件路径
#     json_path = "output11.json"  # 替换为你想要输出的JSON文件路径
#     ppt_to_json(ppt_path, json_path)
    folder_path='D:\服务外包大赛\O_O\data\课件'
    out_path='D:\服务外包大赛\O_O\data\课件(json)'
    pptf_to_json(folder_path,out_path)
